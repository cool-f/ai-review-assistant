"""
TextExtractor — 从课件文件中提取纯文本内容

支持格式: PDF (PyMuPDF), PPTX (python-pptx), DOCX (python-docx), TXT/MD (纯文本)
"""

import base64
import json
import re
from pathlib import Path

from review_assistant.infrastructure.ai.client import get_ai_client
from review_assistant.application.courseware.extractor import AIExtractor


# Vision 逐页提取 prompt（OpenAI content blocks 格式调用）
PER_PAGE_VISION_PROMPT = """You are analyzing page {page_num} of a courseware PDF document. This is a scanned page — extract all educational content visible in the image.

Output ONLY a valid JSON object with this structure:
{{
  "knowledge_points": [
    {{
      "name": "知识点名称",
      "summary": "该页上此知识点的核心内容（100-300字，包含公式/定义/定理）",
      "difficulty": "简单 | 中等 | 困难",
      "page_ref": {page_num},
      "examples": [
        {{
          "question": "例题题目原文",
          "answer": "例题答案",
          "explanation": "解题思路（可为空）"
        }}
      ]
    }}
  ]
}}

Rules:
- Only include knowledge points VISIBLE on this single page
- Mathematical formulas: output in LaTeX notation (e.g., $E=mc^2$)
- Diagrams/tables: describe their key information in text
- If no clear knowledge point, return {{"knowledge_points": []}}
- Output ONLY valid JSON, no extra text"""


# Vision 合并 prompt（纯文本）
VISION_MERGE_PROMPT = """You are a knowledge merging assistant. You will receive knowledge points extracted from multiple pages of a PDF. Your task is to merge them into a single clean list.

Rules:
1. Deduplicate: if the same knowledge point appears on multiple pages, keep the first occurrence and merge summaries if they contain different information
2. Preserve: all unique knowledge points must be kept
3. Keep the original output format (name, summary, difficulty, page_ref, examples)

Output ONLY a valid JSON object with this structure:
{{
  "knowledge_points": [
    {{
      "name": "知识点名称",
      "summary": "合并后的核心内容",
      "difficulty": "简单 | 中等 | 困难",
      "page_ref": 页码,
      "examples": [
        {{
          "question": "例题题目原文",
          "answer": "例题答案",
          "explanation": "解题思路"
        }}
      ]
    }}
  ]
}}

Output ONLY valid JSON, no extra text."""


class TextExtractor:
    """从课件文件中提取纯文本"""

    # 支持的文件类型
    SUPPORTED_TYPES = frozenset({"pdf", "pptx", "docx", "txt", "md"})

    # 最大单块字符数
    MAX_CHUNK_SIZE = 2000

    # ── 公开接口 ──────────────────────────────────
    @classmethod
    def extract(cls, file_path: str, file_type: str) -> tuple[str, int | None]:
        """
        从文件提取纯文本

        Args:
            file_path: 文件绝对路径
            file_type: 文件类型 (pdf/pptx/docx/txt/md)

        Returns:
            (extracted_text, page_count_or_None)

        Raises:
            ValueError: 不支持的文件类型
            FileNotFoundError: 文件不存在
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ft = file_type.lower()

        if ft == "pdf":
            return cls._extract_pdf(str(path))
        elif ft == "pptx":
            return cls._extract_pptx(str(path))
        elif ft == "docx":
            return cls._extract_docx(str(path))
        elif ft in ("txt", "md"):
            return cls._extract_plain(str(path)), None
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

    # ── 分块工具 ──────────────────────────────────
    @classmethod
    def chunk_text(cls, text: str) -> list[dict]:
        """
        将文本按规则分块

        规则:
          1. 先按双换行 (\\n\\n) 分割为段落
          2. 单段 ≤ MAX_CHUNK_SIZE 则直接作为一个 chunk
          3. 超长段按句号/问号/感叹号二次分割（保留分隔符在各自句尾）

        Args:
            text: 原始纯文本

        Returns:
            [{"content": str, "chunk_index": int}, ...]
        """
        if not text or not text.strip():
            return []

        paragraphs = text.split("\n\n")
        chunks: list[dict] = []
        chunk_index = 0
        current_page: int | None = None
        page_marker = re.compile(r"^\[\[PAGE:(\d+)\]\]\s*")

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            marker = page_marker.match(para)
            if marker:
                current_page = int(marker.group(1)) or None
                para = page_marker.sub("", para, count=1).strip()
                if not para:
                    continue

            if len(para) <= cls.MAX_CHUNK_SIZE:
                chunks.append({
                    "content": para,
                    "chunk_index": chunk_index,
                    "page_number": current_page,
                })
                chunk_index += 1
            else:
                # 按句子二次分割
                sub_texts = cls._split_by_sentence(para)
                for sub in sub_texts:
                    chunks.append({
                        "content": sub,
                        "chunk_index": chunk_index,
                        "page_number": current_page,
                    })
                    chunk_index += 1

        return chunks

    @classmethod
    def _split_by_sentence(cls, text: str) -> list[str]:
        """
        按中英文句子分隔符切割长段落，保证每段不超过 MAX_CHUNK_SIZE

        分隔符: 。！？ . ! ?
        使用 look-behind 断言让分隔符留在句尾。
        """
        sentences = re.split(r"(?<=[。！？.!?])", text)
        result: list[str] = []
        current = ""

        for sent in sentences:
            if not sent:
                continue
            if len(current) + len(sent) <= cls.MAX_CHUNK_SIZE:
                current += sent
            else:
                if current.strip():
                    result.append(current.strip())
                # 如果单个句子就超过限制，强制截断
                if len(sent) > cls.MAX_CHUNK_SIZE:
                    for i in range(0, len(sent), cls.MAX_CHUNK_SIZE):
                        piece = sent[i:i + cls.MAX_CHUNK_SIZE].strip()
                        if piece:
                            result.append(piece)
                    current = ""
                else:
                    current = sent

        if current.strip():
            result.append(current.strip())

        return result

    # ── PDF 提取 (PyMuPDF) ────────────────────────
    @staticmethod
    def _extract_pdf(file_path: str) -> tuple[str, int]:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        texts: list[str] = []
        try:
            for page_number, page in enumerate(doc, start=1):
                page_text = page.get_text("text")
                if page_text:
                    texts.append(f"[[PAGE:{page_number}]]\n{page_text.strip()}")
            page_count = doc.page_count
        finally:
            doc.close()

        return "\n\n".join(texts), page_count

    # ── PPTX 提取 (python-pptx) ───────────────────
    @staticmethod
    def _extract_pptx(file_path: str) -> tuple[str, int | None]:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text: list[str] = []

        for page_number, slide in enumerate(prs.slides, start=1):
            shapes_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            shapes_text.append(para_text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if row_texts:
                            shapes_text.append(" | ".join(row_texts))
            if shapes_text:
                slides_text.append(
                    f"[[PAGE:{page_number}]]\n" + "\n".join(shapes_text)
                )

        page_count = len(prs.slides)
        return "\n\n".join(slides_text), page_count

    # ── DOCX 提取 (python-docx) ───────────────────
    @staticmethod
    def _extract_docx(file_path: str) -> tuple[str, int | None]:
        from docx import Document

        doc = Document(file_path)
        paragraphs: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # 也提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                row_texts = [
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                ]
                if row_texts:
                    paragraphs.append(" | ".join(row_texts))

        return "\n\n".join(paragraphs), None

    # ── 纯文本提取 (txt / md) ─────────────────────
    @staticmethod
    def _extract_plain(file_path: str) -> str:
        # 尝试多种编码
        for encoding in ("utf-8", "utf-8-sig", "gb2312", "gbk", "latin-1"):
            try:
                return Path(file_path).read_text(encoding=encoding)
            except (UnicodeDecodeError, UnicodeError):
                continue
        # 最后 fallback: 忽略无法解码的字符
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")

    # ── Vision PDF 提取（图片型 PDF 多模态识别）───
    @classmethod
    async def _extract_pdf_via_vision(cls, file_path: str) -> tuple[list[dict], int]:
        """
        使用视觉识别（Vision API）逐页提取 PDF 中的知识点

        逐页渲染为图片，调用 AI Vision API 识别每页的知识点，
        收集所有页的结果后使用纯文本 AI 做去重合并。

        Args:
            file_path: PDF 文件绝对路径

        Returns:
            (all_knowledge_points_list, page_count)

        Raises:
            RuntimeError: Vision 提取失败
        """
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        all_page_results: list[list[dict]] = []
        page_count = doc.page_count

        try:
            for page_num in range(page_count):
                page = doc[page_num]

                # 渲染页面为 PNG 图片（200 DPI）
                pix = page.get_pixmap(dpi=200)
                img_bytes = pix.tobytes("png")
                img_b64 = base64.b64encode(img_bytes).decode("utf-8")

                # 构建含图片的 content blocks 消息
                prompt_text = PER_PAGE_VISION_PROMPT.format(page_num=page_num + 1)
                messages = [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt_text},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_b64}"
                                },
                            },
                        ],
                    }
                ]

                # 调用 AI
                client = get_ai_client()
                response = await client.chat(
                    messages=messages,
                    temperature=0.3,
                    max_tokens=4096,
                )

                raw_text = response.content.strip()
                page_result = AIExtractor._parse_json_response(raw_text)
                page_kps = page_result.get("knowledge_points", [])
                all_page_results.append(page_kps)

        finally:
            doc.close()

        # 合并所有页的结果
        merged = await cls._merge_vision_results(all_page_results)
        return merged, page_count

    @classmethod
    async def _merge_vision_results(
        cls, all_page_results: list[list[dict]]
    ) -> list[dict]:
        """
        将所有页的 Vision 提取结果合并、去重

        构建合并 prompt，调一次纯文本 AI（无图片），做去重合并。

        Args:
            all_page_results: 每页的知识点列表，每页为 [{"name":..., ...}, ...]

        Returns:
            合并去重后的标准 knowledge_points 列表
        """
        # 展开所有页知识点
        all_kps: list[dict] = []
        for page_kps in all_page_results:
            all_kps.extend(page_kps)

        if not all_kps:
            return []

        # 如果只有一页的结果，直接返回（无需合并去重）
        if len(all_page_results) <= 1:
            return all_kps

        # 序列化当前结果用于合并 prompt
        serialized = json.dumps(all_kps, ensure_ascii=False, indent=2)
        user_text = (
            f"Please merge and deduplicate the following knowledge points "
            f"extracted from a multi-page PDF:\n\n{serialized}"
        )

        messages = [
            {"role": "user", "content": VISION_MERGE_PROMPT},
            {"role": "user", "content": user_text},
        ]

        client = get_ai_client()
        response = await client.chat(
            messages=messages,
            temperature=0.3,
            max_tokens=4096,
        )

        raw_text = response.content.strip()
        merged_result = AIExtractor._parse_json_response(raw_text)
        merged_kps = merged_result.get("knowledge_points", [])

        # 如果合并失败回退到展开原始结果
        if not merged_kps:
            return all_kps

        return merged_kps
